import os
import os.path
import ctypes
import shutil
import time
import pickle
import sys
from string import punctuation
from PyQt5.QtWidgets import QApplication, QWidget, QVBoxLayout, QPushButton, QLabel, QDialog, QMainWindow, QFileDialog, QTextEdit
from deep_translator import GoogleTranslator
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import win32com.client
import pythoncom
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from PyQt5.QtWebEngineWidgets import QWebEngineView
from PyQt5.QtCore import QUrl, pyqtSlot, QSize
from PyQt5.QtGui import QFont
import tempfile
from googleapiclient.http import MediaFileUpload, MediaIoBaseDownload
import requests
from googleapiclient.errors import HttpError
import io
import urllib.parse
from googleapiclient import errors

# If modifying these SCOPES, delete the file token.pickle.
SCOPES = ['https://www.googleapis.com/auth/drive.metadata.readonly']

# Create a translator object
translator = GoogleTranslator(source='auto', target='english')

# Create a Qt application
app = QApplication([])

def authenticate_google_drive():
    creds = None
    SCOPES = ['https://www.googleapis.com/auth/drive.metadata.readonly', 
              'https://www.googleapis.com/auth/userinfo.email', 
              'https://www.googleapis.com/auth/userinfo.profile',
              'openid']

    # Check if token.pickle file exists and is not empty
    if os.path.exists('token.pickle') and os.path.getsize('token.pickle') > 0:
        with open('token.pickle', 'rb') as token:
            creds = pickle.load(token)
    else:
        # If there are no (valid) credentials available, let the user log in.
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(
                'credentials.json', SCOPES)  # here enter the name of your downloaded JSON file
            creds = flow.run_local_server(port=0)
        # Save the credentials for the next run
        with open('token.pickle', 'wb') as token:
            pickle.dump(creds, token)
            
    # Call the Drive v3 API
    drive_service = build('drive', 'v3', credentials=creds)
    # Call the People API to get the user's email
    people_service = build('people', 'v1', credentials=creds)
    profile = people_service.people().get(resourceName='people/me', personFields='emailAddresses').execute()
    email = profile['emailAddresses'][0]['value']
    
    return drive_service, email

def download_file(drive_service, file_id, directory_path, file_name, translated_root_id, override=False):
    if not os.path.exists(directory_path):
        os.makedirs(directory_path)  # Create the directory if it doesn't exist

    translated_file_name = translate_file_name(file_name)
    file_path = os.path.join(directory_path, translated_file_name)

    # Check if the file already exists in the translated_root_id directory on Google Drive
    encoded_filename = urllib.parse.quote(translated_file_name)
    response = drive_service.files().list(
        q=f"name='{encoded_filename}' and '{translated_root_id}' in parents",
        fields='files(id, name)').execute()
    
    if response.get('files') and not override:
        print(f"File '{translated_file_name}' already exists in parent folder ID '{translated_root_id}'")
        return None  # Return None if the file already exists
    
    elif response.get('files') and override:
        # Get the ID of the file to delete
        file_to_delete_id = response.get('files')[0]['id']
        delete_file(drive_service, file_to_delete_id)  # Delete the existing file
    
    # If the file does not exist, download it
    request = drive_service.files().get_media(fileId=file_id)
    fh = io.BytesIO()
    downloader = MediaIoBaseDownload(fh, request)
    done = False
    while done is False:
        status, done = downloader.next_chunk()
    with open(file_path, 'wb') as f:
        f.write(fh.getbuffer())

    return file_path

def upload_file(drive_service, file_path, parent_folder_id, override=True):
    file_metadata = {
        'name': os.path.basename(file_path),
        'parents': [parent_folder_id]
    }

    # Check if the file already exists in the parent_folder_id directory on Google Drive
    response = drive_service.files().list(
        q=f"name='{file_metadata['name']}' and '{parent_folder_id}' in parents",
        fields='files(id, name)').execute()
    
    if response.get('files') and not override:
        print(f"File '{file_metadata['name']}' already exists in parent folder ID '{parent_folder_id}'")
        return  # Do not upload the file if it already exists
    
    elif response.get('files') and override:
        # Get the ID of the file to delete
        file_to_delete_id = response.get('files')[0]['id']
        delete_file(drive_service, file_to_delete_id)  # Delete the existing file
    
    media = MediaFileUpload(file_path)
    file = drive_service.files().create(body=file_metadata, media_body=media, fields='id').execute()
    print(f"File uploaded with ID: {file.get('id')}")

def delete_file(drive_service, file_id):
    try:
        drive_service.files().delete(fileId=file_id).execute()
    except errors.HttpError as error:
        print(f'An error occurred: {error}')

def select_directory():
    dialog = QFileDialog()
    dialog.setFileMode(QFileDialog.Directory)
    dialog.exec_()
    return dialog.selectedFiles()[0]

def is_hidden(filepath):
    name = os.path.basename(os.path.abspath(filepath))
    return name.startswith('.') or has_hidden_attribute(filepath)

def has_hidden_attribute(filepath):
    try:
        attribute = ctypes.windll.kernel32.GetFileAttributesW(str(filepath))
        assert attribute != -1
        result = attribute & 2
    except (AttributeError, AssertionError):
        result = False
    return result

def log_message(log, message):
    # Add message to log
    log.append(message)  # QTextEdit.append automatically adds a newline

    # Process events to make sure GUI updates immediately
    QApplication.processEvents()

def upload_directory_to_drive(drive_service, directory_path, parent_folder_id):
    for file_name in os.listdir(directory_path):
        file_path = os.path.join(directory_path, file_name)

        if os.path.isfile(file_path):
            # Upload file
            file_metadata = {
                'name': file_name,
                'parents': [parent_folder_id]
            }
            media = MediaFileUpload(file_path)
            drive_service.files().create(body=file_metadata, media_body=media, fields='id').execute()
        else:
            # Create folder on Drive
            folder_metadata = {
                'name': file_name,
                'mimeType': 'application/vnd.google-apps.folder',
                'parents': [parent_folder_id]
            }
            folder = drive_service.files().create(body=folder_metadata, fields='id').execute()

            # Upload files in the folder
            upload_directory_to_drive(drive_service, file_path, folder.get('id'))

def download_directory_from_drive(drive_service, folder_id, local_directory):
    query = f"'{folder_id}' in parents and trashed=false"
    results = drive_service.files().list(q=query).execute()
    items = results.get('files', [])

    if not items:
        print('No files found.')
    else:
        for item in items:
            # If item is a file
            if item['mimeType'] != 'application/vnd.google-apps.folder':
                request = drive_service.files().get_media(fileId=item['id'])
                fh = io.FileIO(os.path.join(local_directory, item['name']), 'wb')
                downloader = MediaIoBaseDownload(fh, request)
                done = False
                while done is False:
                    status, done = downloader.next_chunk()
                print(f"Downloaded file: {item['name']}")
            # If item is a directory
            else:
                sub_dir_path = os.path.join(local_directory, item['name'])
                if not os.path.exists(sub_dir_path):
                    os.makedirs(sub_dir_path)
                # Recursively download the subdirectory
                download_directory_from_drive(drive_service, item['id'], sub_dir_path)

def translate_files_local_to_local(source_directory, target_directory, label, log):
    # Define a set of file extensions for which the name translation should be skipped
    skip_translation_extensions = {'.mp4', '.wav', '.avi', '.mov', '.mpg', '.wmv'}

    for item in os.listdir(source_directory):
        source_item_path = os.path.join(source_directory, item)

        # Check if the current file extension is in the skip_translation_extensions set
        if os.path.splitext(item)[1].lower() in skip_translation_extensions:
            target_item_path = os.path.join(target_directory, item)

            # Create a text file with the translated name of the video file
            translated_video_name = translate_file_name(item)
            video_translation_file_path = os.path.join(target_directory, translated_video_name + '.txt')
            with open(video_translation_file_path, 'w', encoding='utf-8') as file:
                file.write(f"{item} = {translated_video_name}")
        else:
            target_item_path = os.path.join(target_directory, translate_file_name(item))

        # Skip hidden files and directories
        if is_hidden(source_item_path):
            continue

        if os.path.isdir(source_item_path):
            os.makedirs(target_item_path, exist_ok=True)
            translate_files_local_to_local(source_item_path, target_item_path, label, log)  # Don't forget to pass the log
        else:
            log_message(log, f'Translating {item}...')  # Replaces `label.setText(f'Translating {item}...')`

            if item.endswith(('.doc', '.docx', '.pptx')):
                try:
                    shutil.copy(source_item_path, target_item_path)
                except Exception as e:
                    log_message(log, f"Failed to copy file {source_item_path} to {target_item_path}: {e}")  # Replaces `print(...)`
                    continue

                if item.endswith(('.docx', '.doc')):
                    if item.endswith('.doc'):  # If it is a .doc file
                        temp = convert_doc_to_docx(target_item_path)
                        os.remove(target_item_path)
                        target_item_path = temp
                    translate_docx(target_item_path)
                elif item.endswith('.pptx'):
                    translate_pptx(target_item_path)
            else:
                try:
                    shutil.copy(source_item_path, target_item_path)
                except Exception as e:
                    log_message(log, f"Failed to copy file {source_item_path} to {target_item_path}: {e}")  # Replaces `print(...)`
    log_message(log, 'Translation finished')  # Replaces `label.setText('Translation finished')`

def translate_files_local_to_drive(source_directory, target_folder_id, drive_service, label, log):
    # Create a temporary directory to store translated files
    with tempfile.TemporaryDirectory() as temp_dir:
        # Translate files
        translate_files_local_to_local(source_directory, temp_dir, label, log)
        # Upload translated files to Drive
        upload_directory_to_drive(drive_service, temp_dir, target_folder_id)

def translate_files_drive_to_local(source_folder_id, target_directory, drive_service, label, log):
    # Create a temporary directory to store the downloaded files
    with tempfile.TemporaryDirectory() as temp_dir:
        # Download files from Drive
        download_directory_from_drive(drive_service, source_folder_id, temp_dir)
        # Translate downloaded files and save them to target_directory
        translate_files_local_to_local(temp_dir, target_directory, label, log)

def create_folder(drive_service, name, parent_id):
    name = sanitize_name(name)  # Remove illegal characters from the folder name

    # Check if the folder already exists
    response = drive_service.files().list(
        q=f"name='{name}' and '{parent_id}' in parents and mimeType='application/vnd.google-apps.folder'",
        spaces='drive',
        fields='files(id, name)').execute()
    
    if response.get('files'):
        folder_id = response.get('files')[0].get('id')
        print(f"Folder '{name}' already exists in parent folder ID '{parent_id}'")
        return folder_id, False  # Return False indicating the folder already exists

    # Create the folder if it doesn't exist
    file_metadata = {
        'name': name,
        'mimeType': 'application/vnd.google-apps.folder',
        'parents': [parent_id]
    }
    file = drive_service.files().create(body=file_metadata,
                                        fields='id').execute()
    
    print(f"Folder '{name}' created in parent folder ID '{parent_id}'")
    return file.get('id'), True  # Return True indicating the folder is newly created

def copy_and_rename_file(drive_service, file_id, translated_root_id, translated_file_name, override=False):
    # Make a copy of the original file in the new directory
    file_metadata = {
        'name': translated_file_name,
        'parents': [translated_root_id]
    }

    # Check if the file already exists in the translated_root_id directory on Google Drive
    response = drive_service.files().list(
        q=f"name='{translated_file_name}' and '{translated_root_id}' in parents",
        fields='files(id, name)').execute()
    
    if response.get('files') and not override:
        print(f"File '{translated_file_name}' already exists in parent folder ID '{translated_root_id}'")
        return  # Do not copy the file if it already exists
    
    elif response.get('files') and override:
        # Get the ID of the file to delete
        file_to_delete_id = response.get('files')[0]['id']
        delete_file(file_to_delete_id)  # Delete the existing file
    
    copied_file = drive_service.files().copy(
        fileId=file_id,
        body=file_metadata,
        fields='id'
    ).execute()
    print(f"File copied and renamed with ID: {copied_file.get('id')}")

def translate_files_drive_to_drive(source_folder_id, target_folder_id, drive_service, label, log, start_file=None, convert_docs=False, override_docs=False, convert_slides=False, override_slides=False, copy_translate_others=False, override_others=False):
    results = drive_service.files().list(
        q=f"'{source_folder_id}' in parents and mimeType='application/vnd.google-apps.folder'",
        fields="files(id, name)").execute()

    items = results.get('files', [])

    for item in items:
        subdirectory_id = item['id']
        file_name = translate_text(item['name'])

        global start_translating
        if start_file and not start_translating and file_name == start_file:
            start_translating = True

        translated_subdirectory_id, is_new_folder = create_folder(drive_service, file_name, target_folder_id)
        if is_new_folder:  # Only process the subdirectory if it is newly created
            print(f"Processing subdirectory '{item['name']}' with ID '{subdirectory_id}'")
            translate_files_drive_to_drive(subdirectory_id, translated_subdirectory_id, drive_service, label, log, start_file, convert_docs, override_docs, convert_slides, override_slides, copy_translate_others, override_others)
        #else:
        #    print(f"Processing subdirectory '{item['name']}' with ID '{subdirectory_id}'")
        #    translate_files_drive_to_drive(subdirectory_id, translated_subdirectory_id, drive_service, label, log, start_file, convert_docs, override_docs, convert_slides, override_slides, copy_translate_others, override_others)

    results = drive_service.files().list(
        q=f"'{source_folder_id}' in parents and mimeType!='application/vnd.google-apps.folder'",
        fields="files(id, name, mimeType)").execute()
    items = results.get('files', [])

    local_directory_path = os.path.join('./temp_drive_files', source_folder_id)
    for item in items:
        if not start_translating:
            continue

        file_name = item['name']
        file_id = item['id']
        translated_file_name = translate_file_name(file_name)

        encoded_filename = urllib.parse.quote(translated_file_name)
        response = drive_service.files().list(
            q=f"name='{encoded_filename}' and '{target_folder_id}' in parents",
            fields='files(id, name)').execute()
        if response.get('files'):
            print(f"File '{translated_file_name}' already exists in parent folder ID '{target_folder_id}'")
            continue

        if item['mimeType'] in ['application/msword', 'application/vnd.openxmlformats-officedocument.wordprocessingml.document'] and convert_docs:
            file_path = download_file(file_id, local_directory_path, file_name, target_folder_id, False)
            if item['mimeType'] == 'application/msword':
                file_path = convert_doc_to_docx(file_path)
            translated_file_path = translate_docx(file_path)
            upload_file(translated_file_path, target_folder_id, override_docs)
            os.remove(translated_file_path)
        elif item['mimeType'] == 'application/vnd.openxmlformats-officedocument.presentationml.presentation' and convert_slides:
            file_path = download_file(file_id, local_directory_path, file_name, target_folder_id, False)
            translated_file_path = translate_pptx(file_path)
            upload_file(translated_file_path, target_folder_id, override_slides)
            os.remove(translated_file_path)
        elif copy_translate_others:
            translated_file_name = translate_file_name(file_name)
            copy_and_rename_file(file_id, target_folder_id, translated_file_name, override_others)

def sanitize_name(name):
    invalid_chars = ['<', '>', ':', '"', '/', '\\', '|', '?', '*']
    for char in invalid_chars:
        name = name.replace(char, '')
    return name

def translate_file_name(file):
    translated_name = translator.translate(file)
    return sanitize_name(translated_name)

def is_punctuation(text):
    return all(char in punctuation for char in text)

def split_text_into_chunks(text, chunk_size=5000):
    """
    Function to split the text into smaller chunks
    """
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

def translate_text(text):
    """
    Function to translate the text
    """
    text = str(text)  # Ensure the text is a string

    # Return the original text if it has less than 2 non-whitespace characters 
    # or consists only of punctuation or if it is a digit
    if len(text.strip()) < 2 or is_punctuation(text.strip()) or text.isdigit():
        return text

    chunks = split_text_into_chunks(text)  # Split the text into chunks
    translated_text = ""  # Placeholder for the translated text

    for chunk in chunks:
        while True:
            try:
                translated_chunk = translator.translate(chunk)  # Translate the chunk
                if translated_chunk is None:
                    print(f"Failed to translate text-english: {chunk}. Result was None.")
                else:
                    translated_text += translated_chunk  # Add the translated chunk to the translated text
                break  # Break the while loop if the translation was successful

            except Exception as e:
                print(f"Failed to translate text-english: {chunk}. Error: {str(e)}")
                time.sleep(1)  # Wait for 1 second before retrying

    #translated_text = translated_text.replace("'", "")  # remove apostrophes
    return translated_text

def translate_docx(translated_file):
    doc = Document(translated_file)
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue  # Skip empty runs
            try:
                translated_text = translate_text(run.text)
                run.text = translated_text
            except Exception as e:
                print(f"Failed to translate text: {run.text}. Error: {str(e)}")
                continue
    
    # Translate text in tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if not run.text:
                            continue  # Skip empty runs
                        try:
                            translated = translate_text(run.text)  # Translate the text to English
                            run.text = translated  # Update the text with the translated version
                        except Exception as e:
                            print(f"Failed to translate text: {run.text}. Error: {str(e)}")
                            continue

    directory, filename = os.path.split(translated_file)
    filename_without_ext = os.path.splitext(filename)[0]

    try:
        translated_filename_without_ext = translate_text(filename_without_ext)
    except Exception as e:
        print(f"Failed to translate filename: {filename_without_ext}. Error: {str(e)}")
        translated_filename_without_ext = filename_without_ext

    translated_filename = translated_filename_without_ext + ".docx"
    translated_doc_path = os.path.join(directory, translated_filename)
    doc.save(translated_doc_path)

    print(f"Translated document saved at: {translated_doc_path}")
    return translated_doc_path

def translate_text_frame(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue  # Skip empty runs
            try:
                translated = translate_text(run.text)  # Translate the text to English
                run.text = translated  # Update the text with the translated version
            except Exception as e:
                print(f"Failed to translate text-frame: {run.text}. Error: {str(e)}")
                continue

def adjust_text_size(shape):
    while shape.text_frame.text != "":
        try:
            # Try to access the last character of the shape's text
            _ = shape.text_frame.text[-1]
            break  # Break the loop if the last character can be accessed
        except IndexError:
            # If the last character cannot be accessed, the text overflows the shape
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Decrease the font size by 1 point
                    run.font.size = Pt(run.font.size.pt - 1)

def process_shape(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP: # type: ignore
        # If the shape is a group, recursively process each shape within the group
        for shape_in_group in shape.shapes:
            process_shape(shape_in_group)
    elif shape.has_text_frame:
        # If the shape has a text frame, translate the text and adjust the text size
        translate_text_frame(shape.text_frame)
        adjust_text_size(shape)
    elif shape.has_table:
        # If the shape is a table, process each cell's text frame and adjust text size
        for row in shape.table.rows:
            for cell in row.cells:
                translate_text_frame(cell.text_frame)
                adjust_text_size(cell)

def translate_pptx(pptx_path):
    pres = Presentation(pptx_path)

    # Process each slide in the presentation
    for slide in pres.slides:
        for shape in slide.shapes:
            process_shape(shape)
            if shape.has_text_frame:
                translate_text_frame(shape.text_frame)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        translate_text_frame(cell.text_frame)

    directory, filename = os.path.split(pptx_path)
    filename_without_ext = os.path.splitext(filename)[0]

    try:
        translated_filename_without_ext = translate_text(filename_without_ext)
    except Exception as e:
        print(f"Failed to translate filename: {filename_without_ext}. Error: {str(e)}")
        translated_filename_without_ext = filename_without_ext

    translated_filename = translated_filename_without_ext + ".pptx"
    translated_pptx_path = os.path.join(directory, translated_filename)
    pres.save(translated_pptx_path)

    print(f"Translated presentation saved at: {translated_pptx_path}")
    return translated_pptx_path

def convert_doc_to_docx(doc_path):
    # Ensure the path is absolute
    doc_path = os.path.abspath(doc_path)

    # Create the new path by replacing the extension
    new_file_abs = doc_path.replace(".doc", ".docx")

    try:
        # Initialize the Word.Application
        word = win32com.client.Dispatch('Word.Application',pythoncom.CoInitialize())

        # Set the application to be invisible
        word.Visible = False

        # Open the document
        doc = word.Documents.Open(doc_path)

        # Save as a .docx file
        doc.SaveAs(new_file_abs, FileFormat=16)  # 16 represents the wdFormatDocx constant

        # Close the document
        doc.Close()

        # Quit Word
        word.Quit()
    except Exception as e:
        print(f"Failed to convert file: {doc_path}. Error: {str(e)}")
        return doc_path

    return new_file_abs

def start_translations(source_directory_local, source_directory_drive, target_directory_local, target_directory_drive, label, drive, log):
    if source_directory_local and target_directory_local:
        translate_files_local_to_local(source_directory_local, target_directory_local, label, log)
    elif source_directory_local and target_directory_drive:
        translate_files_local_to_drive(source_directory_local, target_directory_drive, label, drive, log)
    elif source_directory_drive and target_directory_local:
        translate_files_drive_to_local(source_directory_drive, target_directory_local, label, drive, log)
    elif source_directory_drive and target_directory_drive:
        translate_files_drive_to_drive(source_directory_drive, target_directory_drive, label, drive, log)

class MainWindow(QMainWindow):
    def __init__(self):
        super(MainWindow, self).__init__()

        # Initialize source_directory_drive and target_directory_drive
        self.source_directory_drive = None
        self.source_directory_local = None
        self.target_directory_drive = None
        self.target_directory_local = None

        # Checks
        self.isSource = [False]
        self.isTarget = [False]
        self.isDrive = [False]

        self.setWindowTitle('My Window')

        # Create QFont object for buttons and labels
        self.button_font = QFont()
        self.button_font.setPointSize(14)  # Change this value to adjust the font size

        google_drive_button = QPushButton('Authenticate with Google Drive')
        google_drive_button.setFont(self.button_font)
        source_button_local = QPushButton('Select Source Directory (Local)')
        source_button_local.setFont(self.button_font)
        source_button_drive = QPushButton('Select Source Directory (Google Drive)')
        source_button_drive.setFont(self.button_font)
        target_button_local = QPushButton('Select Target Directory (Local)')
        target_button_local.setFont(self.button_font)
        target_button_drive = QPushButton('Select Target Directory (Google Drive)')
        target_button_drive.setFont(self.button_font)
        self.start_button = QPushButton('Start Translations')
        self.start_button.setFont(self.button_font)

        # Create QFont object for labels
        label_font = QFont()
        label_font.setPointSize(12)  # Change this value to adjust the font size

        self.google_user_label = QLabel("Google User: Not Authenticated")
        self.google_user_label.setFont(label_font)
        label = QLabel()
        label.setFont(label_font)
        self.source_directory_label_local = QLabel("Source Directory (Local): Not Selected")
        self.source_directory_label_local.setFont(label_font)
        self.source_directory_label_drive = QLabel("Source Directory (Drive): Not Selected")
        self.source_directory_label_drive.setFont(label_font)
        self.target_directory_label_local = QLabel("Target Directory (Local): Not Selected")
        self.target_directory_label_local.setFont(label_font)
        self.target_directory_label_drive = QLabel("Target Directory (Drive): Not Selected")
        self.target_directory_label_drive.setFont(label_font)
        # Initialize QTextEdit for status log
        self.log = QTextEdit()
        self.log.setReadOnly(True)

        google_drive_button.clicked.connect(
            lambda: (
                setattr(self, 'drive', authenticate_google_drive()[0]),
                self.google_user_label.setText(f'Google User: {authenticate_google_drive()[1]}'),
                source_button_drive.setEnabled(True),
                self.isDrive.__setitem__(0, True)
            )
        )
        source_button_local.clicked.connect(
            lambda: (
                self.clear_source_directory(),  # Clear existing source directory
                setattr(self, 'source_directory_local', select_directory()),
                self.source_directory_label_local.setText(f'Source Directory (Local): {self.source_directory_local}'),
                target_button_local.setEnabled(True),
                target_button_drive.setEnabled(True and self.isDrive[0]),
                self.isSource.__setitem__(0, True)
            )
        )
        source_button_drive.clicked.connect(
            lambda: (
                self.clear_source_directory(),  # Clear existing source directory
                self.enter_folder_id(self.source_directory_label_drive, is_source=True),  # Pass 'is_source' parameter
                target_button_local.setEnabled(True),
                target_button_drive.setEnabled(True and self.isDrive[0]),
                self.isSource.__setitem__(0, True)
            )
        )
        target_button_local.clicked.connect(
            lambda: (
                self.clear_target_directory(),  # Clear existing target directory
                setattr(self, 'target_directory_local', select_directory()),
                self.target_directory_label_local.setText(f'Target Directory (Local): {self.target_directory_local}'),
                self.start_button.setEnabled(True),
                self.isTarget.__setitem__(0, True)
            )
        )
        target_button_drive.clicked.connect(
            lambda: (
                self.clear_target_directory(),  # Clear existing target directory
                self.enter_folder_id(self.target_directory_label_drive, is_source=False),  # Pass 'is_source' parameter
                self.start_button.setEnabled(True),
                self.isTarget.__setitem__(0, True)
            )
        )
        self.start_button.clicked.connect(
            lambda: start_translations(
                self.source_directory_local,
                self.source_directory_drive,
                self.target_directory_local,
                self.target_directory_drive,
                label,
                getattr(self, 'drive', None),
                self.log
            )
        )

        source_button_drive.setEnabled(False)
        target_button_local.setEnabled(False)
        target_button_drive.setEnabled(False)
        self.start_button.setEnabled(False)

        layout = QVBoxLayout()
        layout.addWidget(google_drive_button)
        layout.addWidget(self.google_user_label)
        layout.addWidget(source_button_local)
        layout.addWidget(source_button_drive)
        layout.addWidget(target_button_local)
        layout.addWidget(target_button_drive)
        layout.addWidget(self.source_directory_label_local)
        layout.addWidget(self.source_directory_label_drive)
        layout.addWidget(self.target_directory_label_local)
        layout.addWidget(self.target_directory_label_drive)
        layout.addWidget(self.start_button)
        layout.addWidget(label)
        layout.addWidget(self.log)  # Add status log to layout

        central_widget = QWidget()
        central_widget.setLayout(layout)
        self.setCentralWidget(central_widget)


    def log_message(self, message):
        # Update status label
        self.label.setText(message)

        # Add message to log
        self.log.append(message)  # QTextEdit.append automatically adds a newline

        # Process events to make sure GUI updates immediately
        QApplication.processEvents()

    def clear_source_directory(self):
        # Reset both local and drive source directories
        self.source_directory_local = None
        self.source_directory_drive = None
        # Also reset the labels
        self.source_directory_label_local.setText("Source Directory (Local): Not Selected")
        self.source_directory_label_drive.setText("Source Directory (Drive): Not Selected")

    def clear_target_directory(self):
        # Reset both local and drive target directories
        self.target_directory_local = None
        self.target_directory_drive = None
        # Also reset the labels
        self.target_directory_label_local.setText("Target Directory (Local): Not Selected")
        self.target_directory_label_drive.setText("Target Directory (Drive): Not Selected")

    def enter_folder_id(self, label, is_source):
        # Create a new web view widget
        view = QWebEngineView()

        # Load the Google Drive Picker URL
        view.load(QUrl('https://drive.google.com/drive/folders'))

        # Show the web view widget in a new window
        dialog = QDialog(self)
        dialog.resize(800, 600)  # Change these values to adjust the dialog size
        layout = QVBoxLayout()

        choose_directory_button = QPushButton("Choose Current Directory")
        choose_directory_button.setFont(self.button_font)  # Make sure to use the same font for this button
        choose_directory_button.clicked.connect(lambda: self.choose_directory(view, label, is_source))  # Pass 'is_source' parameter
        layout.addWidget(view)
        layout.addWidget(choose_directory_button)

        dialog.setLayout(layout)
        dialog.exec_()

    def choose_directory(self, view, label, is_source):
        # Get the selected folder ID from the URL
        url = view.url().toString()
        folder_id = url.split('/')[-1]

        # Set the folder ID to the corresponding label
        label.setText(f"Selected Directory (Drive): {folder_id}")

        # Save the folder ID to the appropriate attribute in the MainWindow object
        if is_source:
            self.source_directory_drive = folder_id
        else:
            self.target_directory_drive = folder_id

        # Enable the start button if both source and target directories are selected
        if self.source_directory_drive and self.target_directory_drive:
            self.start_button.setEnabled(True)
        # Add these new methods
    
def main():
    try:
        app = QApplication(sys.argv)
        mainWin = MainWindow()
        mainWin.resize(800, 600)  # Change these values to adjust the main window size
        mainWin.show()
        sys.exit(app.exec_())
    except Exception as e:
        print(f"Exception occurred: {e}")
        # re-raise the exception after printing
        raise

if __name__ == "__main__":
    main()