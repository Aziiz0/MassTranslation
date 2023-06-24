import os
import ctypes
import shutil
import time
import pickle
import sys
from string import punctuation
from PyQt5.QtWidgets import QApplication, QWidget, QVBoxLayout, QPushButton, QLabel, QDialog, QMainWindow, QFileDialog, QTextEdit
from deep_translator import GoogleTranslator
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import win32com.client
import pythoncom
from google_auth_oauthlib.flow import InstalledAppFlow
from google.auth.transport.requests import Request
from googleapiclient.discovery import build
from PyQt5.QtWebEngineWidgets import QWebEngineView
from PyQt5.QtCore import QUrl

# Create a translator object
translator = GoogleTranslator(source='auto', target='english')

# Create a Qt application
app = QApplication([])

def authenticate_google_drive():
    creds = None
    SCOPES = ['https://www.googleapis.com/auth/drive.metadata.readonly', 
              'https://www.googleapis.com/auth/userinfo.email', 
              'https://www.googleapis.com/auth/userinfo.profile',
              'openid']

    # Check if token.pickle file exists and is not empty
    if os.path.exists('token.pickle') and os.path.getsize('token.pickle') > 0:
        with open('token.pickle', 'rb') as token:
            creds = pickle.load(token)
    else:
        # If there are no (valid) credentials available, let the user log in.
        if creds and creds.expired and creds.refresh_token:
            creds.refresh(Request())
        else:
            flow = InstalledAppFlow.from_client_secrets_file(
                'credentials.json', SCOPES)  # here enter the name of your downloaded JSON file
            creds = flow.run_local_server(port=0)
        # Save the credentials for the next run
        with open('token.pickle', 'wb') as token:
            pickle.dump(creds, token)
            
    # Call the Drive v3 API
    drive_service = build('drive', 'v3', credentials=creds)
    # Call the People API to get the user's email
    people_service = build('people', 'v1', credentials=creds)
    profile = people_service.people().get(resourceName='people/me', personFields='emailAddresses').execute()
    email = profile['emailAddresses'][0]['value']
    
    return drive_service, email

def upload_file_to_google_drive(drive, local_file_path, google_folder_id):
    gfile = drive.CreateFile({"title": os.path.basename(local_file_path), "parents": [{"id": google_folder_id}]})
    gfile.SetContentFile(local_file_path)
    gfile.Upload()

def download_file_from_google_drive(drive, google_file_id, local_directory):
    gfile = drive.CreateFile({'id': google_file_id})
    local_file_path = os.path.join(local_directory, gfile['title'])
    gfile.GetContentFile(local_file_path)
    return local_file_path

def get_google_drive_folder_file_ids(drive, folder_id):
    file_list = drive.ListFile({'q': f"'{folder_id}' in parents"}).GetList()
    return [file['id'] for file in file_list]

def select_directory():
    dialog = QFileDialog()
    dialog.setFileMode(QFileDialog.Directory)
    dialog.exec_()
    return dialog.selectedFiles()[0]

def is_hidden(filepath):
    name = os.path.basename(os.path.abspath(filepath))
    return name.startswith('.') or has_hidden_attribute(filepath)

def has_hidden_attribute(filepath):
    try:
        attribute = ctypes.windll.kernel32.GetFileAttributesW(str(filepath))
        assert attribute != -1
        result = attribute & 2
    except (AttributeError, AssertionError):
        result = False
    return result

def log_message(log, message):
    # Add message to log
    log.append(message)  # QTextEdit.append automatically adds a newline

    # Process events to make sure GUI updates immediately
    QApplication.processEvents()

def translate_files_local_to_local(source_directory, target_directory, label, log):
    # Define a set of file extensions for which the name translation should be skipped
    skip_translation_extensions = {'.mp4', '.wav', '.avi', '.mov', '.mpg', '.wmv'}

    for item in os.listdir(source_directory):
        source_item_path = os.path.join(source_directory, item)

        # Check if the current file extension is in the skip_translation_extensions set
        if os.path.splitext(item)[1].lower() in skip_translation_extensions:
            target_item_path = os.path.join(target_directory, item)

            # Create a text file with the translated name of the video file
            translated_video_name = translate_file_name(item)
            video_translation_file_path = os.path.join(target_directory, translated_video_name + '.txt')
            with open(video_translation_file_path, 'w', encoding='utf-8') as file:
                file.write(f"{item} = {translated_video_name}")
        else:
            target_item_path = os.path.join(target_directory, translate_file_name(item))

        # Skip hidden files and directories
        if is_hidden(source_item_path):
            continue

        if os.path.isdir(source_item_path):
            os.makedirs(target_item_path, exist_ok=True)
            translate_files_local_to_local(source_item_path, target_item_path, label, log)  # Don't forget to pass the log
        else:
            log_message(log, f'Translating {item}...')  # Replaces `label.setText(f'Translating {item}...')`

            if item.endswith(('.doc', '.docx', '.pptx')):
                try:
                    shutil.copy(source_item_path, target_item_path)
                except Exception as e:
                    log_message(log, f"Failed to copy file {source_item_path} to {target_item_path}: {e}")  # Replaces `print(...)`
                    continue

                if item.endswith(('.docx', '.doc')):
                    if item.endswith('.doc'):  # If it is a .doc file
                        temp = convert_doc_to_docx(target_item_path)
                        os.remove(target_item_path)
                        target_item_path = temp
                    translate_docx(target_item_path)
                elif item.endswith('.pptx'):
                    translate_pptx(target_item_path)
            else:
                try:
                    shutil.copy(source_item_path, target_item_path)
                except Exception as e:
                    log_message(log, f"Failed to copy file {source_item_path} to {target_item_path}: {e}")  # Replaces `print(...)`
    log_message(log, 'Translation finished')  # Replaces `label.setText('Translation finished')`

def translate_files_local_to_drive(source_directory, google_folder_id, drive, label):
    # similar to translate_files_local_to_local, but instead of copying files to the target directory, it uploads them to Google Drive
    pass

def translate_files_drive_to_local(google_folder_id, target_directory, drive, label):
    # similar to translate_files_local_to_local, but instead of reading files from the source directory, it downloads them from Google Drive
    pass

def translate_files_drive_to_drive(source_folder_id, target_folder_id, drive, label):
    # similar to translate_files_local_to_local, but instead of reading files from the source directory and copying them to the target directory, it downloads them from the source Google Drive folder and uploads them to the target Google Drive folder
    pass

def sanitize_name(name):
    invalid_chars = ['<', '>', ':', '"', '/', '\\', '|', '?', '*']
    for char in invalid_chars:
        name = name.replace(char, '')
    return name

def translate_file_name(file):
    translated_name = translator.translate(file)
    return sanitize_name(translated_name)

def is_punctuation(text):
    return all(char in punctuation for char in text)

def split_text_into_chunks(text, chunk_size=5000):
    """
    Function to split the text into smaller chunks
    """
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

def translate_text(text):
    """
    Function to translate the text
    """
    text = str(text)  # Ensure the text is a string

    # Return the original text if it has less than 2 non-whitespace characters 
    # or consists only of punctuation or if it is a digit
    if len(text.strip()) < 2 or is_punctuation(text.strip()) or text.isdigit():
        return text

    chunks = split_text_into_chunks(text)  # Split the text into chunks
    translated_text = ""  # Placeholder for the translated text

    for chunk in chunks:
        while True:
            try:
                translated_chunk = translator.translate(chunk)  # Translate the chunk
                if translated_chunk is None:
                    print(f"Failed to translate text-english: {chunk}. Result was None.")
                else:
                    translated_text += translated_chunk  # Add the translated chunk to the translated text
                break  # Break the while loop if the translation was successful

            except Exception as e:
                print(f"Failed to translate text-english: {chunk}. Error: {str(e)}")
                time.sleep(1)  # Wait for 1 second before retrying

    #translated_text = translated_text.replace("'", "")  # remove apostrophes
    return translated_text

def translate_docx(translated_file):
    doc = Document(translated_file)
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue  # Skip empty runs
            try:
                translated_text = translate_text(run.text)
                run.text = translated_text
            except Exception as e:
                print(f"Failed to translate text: {run.text}. Error: {str(e)}")
                continue
    
    # Translate text in tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if not run.text:
                            continue  # Skip empty runs
                        try:
                            translated = translate_text(run.text)  # Translate the text to English
                            run.text = translated  # Update the text with the translated version
                        except Exception as e:
                            print(f"Failed to translate text: {run.text}. Error: {str(e)}")
                            continue

    directory, filename = os.path.split(translated_file)
    filename_without_ext = os.path.splitext(filename)[0]

    try:
        translated_filename_without_ext = translate_text(filename_without_ext)
    except Exception as e:
        print(f"Failed to translate filename: {filename_without_ext}. Error: {str(e)}")
        translated_filename_without_ext = filename_without_ext

    translated_filename = translated_filename_without_ext + ".docx"
    translated_doc_path = os.path.join(directory, translated_filename)
    doc.save(translated_doc_path)

    print(f"Translated document saved at: {translated_doc_path}")
    return translated_doc_path

def translate_text_frame(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue  # Skip empty runs
            try:
                translated = translate_text(run.text)  # Translate the text to English
                run.text = translated  # Update the text with the translated version
            except Exception as e:
                print(f"Failed to translate text-frame: {run.text}. Error: {str(e)}")
                continue

def adjust_text_size(shape):
    while shape.text_frame.text != "":
        try:
            # Try to access the last character of the shape's text
            _ = shape.text_frame.text[-1]
            break  # Break the loop if the last character can be accessed
        except IndexError:
            # If the last character cannot be accessed, the text overflows the shape
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Decrease the font size by 1 point
                    run.font.size = Pt(run.font.size.pt - 1)

def process_shape(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        # If the shape is a group, recursively process each shape within the group
        for shape_in_group in shape.shapes:
            process_shape(shape_in_group)
    elif shape.has_text_frame:
        # If the shape has a text frame, translate the text and adjust the text size
        translate_text_frame(shape.text_frame)
        adjust_text_size(shape)
    elif shape.has_table:
        # If the shape is a table, process each cell's text frame and adjust text size
        for row in shape.table.rows:
            for cell in row.cells:
                translate_text_frame(cell.text_frame)
                adjust_text_size(cell)

def translate_pptx(pptx_path):
    pres = Presentation(pptx_path)

    # Process each slide in the presentation
    for slide in pres.slides:
        for shape in slide.shapes:
            process_shape(shape)
            if shape.has_text_frame:
                translate_text_frame(shape.text_frame)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        translate_text_frame(cell.text_frame)

    directory, filename = os.path.split(pptx_path)
    filename_without_ext = os.path.splitext(filename)[0]

    try:
        translated_filename_without_ext = translate_text(filename_without_ext)
    except Exception as e:
        print(f"Failed to translate filename: {filename_without_ext}. Error: {str(e)}")
        translated_filename_without_ext = filename_without_ext

    translated_filename = translated_filename_without_ext + ".pptx"
    translated_pptx_path = os.path.join(directory, translated_filename)
    pres.save(translated_pptx_path)

    print(f"Translated presentation saved at: {translated_pptx_path}")
    return translated_pptx_path

def convert_doc_to_docx(doc_path):
    # Ensure the path is absolute
    doc_path = os.path.abspath(doc_path)

    # Create the new path by replacing the extension
    new_file_abs = doc_path.replace(".doc", ".docx")

    try:
        # Initialize the Word.Application
        word = win32com.client.Dispatch('Word.Application',pythoncom.CoInitialize())

        # Set the application to be invisible
        word.Visible = False

        # Open the document
        doc = word.Documents.Open(doc_path)

        # Save as a .docx file
        doc.SaveAs(new_file_abs, FileFormat=16)  # 16 represents the wdFormatDocx constant

        # Close the document
        doc.Close()

        # Quit Word
        word.Quit()
    except Exception as e:
        print(f"Failed to convert file: {doc_path}. Error: {str(e)}")
        return doc_path

    return new_file_abs

def start_translations(source_directory_local, source_directory_drive, target_directory_local, target_directory_drive, label, drive, log):
    if source_directory_local and target_directory_local:
        translate_files_local_to_local(source_directory_local, target_directory_local, label, log)
    elif source_directory_local and target_directory_drive:
        translate_files_local_to_drive(source_directory_local, target_directory_drive, label, drive, log)
    elif source_directory_drive and target_directory_local:
        translate_files_drive_to_local(source_directory_drive, target_directory_local, label, drive, log)
    elif source_directory_drive and target_directory_drive:
        translate_files_drive_to_drive(source_directory_drive, target_directory_drive, label, drive, log)

class MainWindow(QMainWindow):
    def __init__(self):
        super(MainWindow, self).__init__()

        # Initialize source_directory_drive and target_directory_drive
        self.source_directory_drive = None
        self.target_directory_drive = None

        self.setWindowTitle('My Window')

        google_drive_button = QPushButton('Authenticate with Google Drive')
        source_button_local = QPushButton('Select Source Directory (Local)')
        source_button_drive = QPushButton('Select Source Directory (Google Drive)')
        target_button_local = QPushButton('Select Target Directory (Local)')
        target_button_drive = QPushButton('Select Target Directory (Google Drive)')
        start_button = QPushButton('Start Translations')
        google_user_label = QLabel("Google User: Not Authenticated")
        label = QLabel()
        source_directory_label_local = QLabel("Source Directory (Local): Not Selected")
        source_directory_label_drive = QLabel("Source Directory (Drive): Not Selected")
        target_directory_label_local = QLabel("Target Directory (Local): Not Selected")
        target_directory_label_drive = QLabel("Target Directory (Drive): Not Selected")

        # Initialize QTextEdit for status log
        self.log = QTextEdit()
        self.log.setReadOnly(True)

        google_drive_button.clicked.connect(
            lambda: (
                setattr(self, 'drive', authenticate_google_drive()[0]),
                google_user_label.setText(f'Google User: {authenticate_google_drive()[1]}'),
                source_button_drive.setEnabled(True)
            )
        )
        source_button_local.clicked.connect(
            lambda: (
                setattr(self, 'source_directory_local', select_directory()),
                source_directory_label_local.setText(f'Source Directory (Local): {self.source_directory_local}'),
                target_button_local.setEnabled(True),
                target_button_drive.setEnabled(True)
            )
        )
        source_button_drive.clicked.connect(
            lambda: (
                self.enter_folder_id(source_directory_label_drive),
                target_button_local.setEnabled(True),
                target_button_drive.setEnabled(True)
            )
        )
        target_button_local.clicked.connect(
            lambda: (
                setattr(self, 'target_directory_local', select_directory()),
                target_directory_label_local.setText(f'Target Directory (Local): {self.target_directory_local}'),
                start_button.setEnabled(True)
            )
        )
        target_button_drive.clicked.connect(
            lambda: (
                self.enter_folder_id(target_directory_label_drive),
                start_button.setEnabled(True)
            )
        )
        start_button.clicked.connect(
            lambda: start_translations(
                self.source_directory_local,
                self.source_directory_drive,
                self.target_directory_local,
                self.target_directory_drive,
                label,
                getattr(self, 'drive', None),
                self.log
            )
        )

        source_button_drive.setEnabled(False)
        target_button_local.setEnabled(False)
        target_button_drive.setEnabled(False)
        start_button.setEnabled(False)

        layout = QVBoxLayout()
        layout.addWidget(google_drive_button)
        layout.addWidget(google_user_label)
        layout.addWidget(source_button_local)
        layout.addWidget(source_button_drive)
        layout.addWidget(target_button_local)
        layout.addWidget(target_button_drive)
        layout.addWidget(source_directory_label_local)
        layout.addWidget(source_directory_label_drive)
        layout.addWidget(target_directory_label_local)
        layout.addWidget(target_directory_label_drive)
        layout.addWidget(start_button)
        layout.addWidget(label)
        layout.addWidget(self.log)  # Add status log to layout

        central_widget = QWidget()
        central_widget.setLayout(layout)
        self.setCentralWidget(central_widget)

    def log_message(self, message):
        # Update status label
        self.label.setText(message)

        # Add message to log
        self.log.append(message)  # QTextEdit.append automatically adds a newline

        # Process events to make sure GUI updates immediately
        QApplication.processEvents()

    def enter_folder_id(self, label):
        # Create a new web view widget
        view = QWebEngineView()

        # Load the Google Drive Picker URL
        view.load(QUrl('https://drive.google.com/drive/folders'))

        # Show the web view widget in a new window
        dialog = QDialog(self)
        layout = QVBoxLayout()

        choose_directory_button = QPushButton("Choose Current Directory")
        choose_directory_button.clicked.connect(lambda: self.choose_directory(view, label))
        layout.addWidget(view)
        layout.addWidget(choose_directory_button)

        dialog.setLayout(layout)
        dialog.exec_()

    def choose_directory(self, view, label):
        # Get the selected folder ID from the URL
        url = view.url().toString()
        folder_id = url.split('/')[-1]

        # Set the folder ID to the corresponding label
        label.setText(f"Selected Directory (Drive): {folder_id}")

        # Save the folder ID to the appropriate attribute in the MainWindow object
        if label.text().startswith("Source"):
            self.source_directory_drive = folder_id
        elif label.text().startswith("Target"):
            self.target_directory_drive = folder_id

        # Enable the start button if both source and target directories are selected
        if self.source_directory_drive and self.target_directory_drive:
            self.start_button.setEnabled(True)

def main():
    try:
        app = QApplication(sys.argv)
        QApplication.setApplicationName('My Window')
        window = MainWindow()

        window.show()

        app.exec_()
    except Exception as e:
        print(f"Exception occurred: {e}")
        # re-raise the exception after printing
        raise

if __name__ == "__main__":
    main()