import os
import shutil
import time
from string import punctuation
from PyQt5.QtWidgets import QApplication, QFileDialog, QVBoxLayout, QWidget, QPushButton, QLabel
from deep_translator import GoogleTranslator
from docx import Document
from pptx import Presentation
from pptx.util import Pt
from pptx.enum.shapes import MSO_SHAPE_TYPE
import win32com.client
import pythoncom

# Create a translator object
translator = GoogleTranslator(source='auto', target='english')

# Create a Qt application
app = QApplication([])

def select_directory():
    dialog = QFileDialog()
    dialog.setFileMode(QFileDialog.Directory)
    dialog.exec_()
    return dialog.selectedFiles()[0]

def translate_files(source_directory, target_directory, label):
    for root, dirs, files in os.walk(source_directory):
        # Translate and create directories first
        for dir in dirs:
            original_dir_path = os.path.join(root, dir)
            translated_dir_path = original_dir_path.replace(source_directory, target_directory)
            translated_dir_name = translate_file_name(os.path.basename(translated_dir_path))
            translated_dir_path = os.path.join(os.path.dirname(translated_dir_path), translated_dir_name)

            if not os.path.exists(translated_dir_path):
                os.makedirs(translated_dir_path)
        
        # Then translate and copy the files
        for file in files:
            original_file_path = os.path.join(root, file)
            translated_file_path = original_file_path.replace(source_directory, target_directory)
            translated_file_name = translate_file_name(os.path.basename(translated_file_path))
            translated_file_path = os.path.join(os.path.dirname(translated_file_path), translated_dir_name, translated_file_name)

            translated_file_dir = os.path.dirname(translated_file_path)
            if not os.path.exists(translated_file_dir):
                os.makedirs(translated_file_dir)

            label.setText(f'Translating {file}...')
            QApplication.processEvents()

            if file.endswith(('.doc', '.docx', '.pptx')):
                try:
                    shutil.copy(original_file_path, translated_file_path)
                except Exception as e:
                    print(f"Failed to copy file {original_file_path} to {translated_file_path}: {e}")
                    continue

                if file.endswith(('.docx', '.doc')):
                    if file.endswith('.doc'):  # If it is a .doc file
                        temp = convert_doc_to_docx(translated_file_path)
                        os.remove(translated_file_path)
                        translated_file_path = temp
                    translate_docx(translated_file_path)
                elif file.endswith('.pptx'):
                    translate_pptx(translated_file_path)
            else:
                try:
                    shutil.copy(original_file_path, translated_file_path)
                except Exception as e:
                    print(f"Failed to copy file {original_file_path} to {translated_file_path}: {e}")

    label.setText('Translation finished')


def sanitize_name(name):
    invalid_chars = ['<', '>', ':', '"', '/', '\\', '|', '?', '*']
    for char in invalid_chars:
        name = name.replace(char, '')
    return name

def translate_file_name(file):
    translated_name = translator.translate(file)
    return sanitize_name(translated_name)

def is_punctuation(text):
    return all(char in punctuation for char in text)

def split_text_into_chunks(text, chunk_size=5000):
    """
    Function to split the text into smaller chunks
    """
    return [text[i:i+chunk_size] for i in range(0, len(text), chunk_size)]

def translate_text(text):
    """
    Function to translate the text
    """
    text = str(text)  # Ensure the text is a string

    # Return the original text if it has less than 2 non-whitespace characters 
    # or consists only of punctuation or if it is a digit
    if len(text.strip()) < 2 or is_punctuation(text.strip()) or text.isdigit():
        return text

    chunks = split_text_into_chunks(text)  # Split the text into chunks
    translated_text = ""  # Placeholder for the translated text

    for chunk in chunks:
        while True:
            try:
                translated_chunk = translator.translate(chunk)  # Translate the chunk
                # print(f"Translated chunk: {translated_chunk}")
                translated_text += translated_chunk  # Add the translated chunk to the translated text
                break  # Break the while loop if the translation was successful
            except Exception as e:
                print(f"Failed to translate text-english: {chunk}. Error: {str(e)}")
                time.sleep(1)  # Wait for 1 second before retrying

    #translated_text = translated_text.replace("'", "")  # remove apostrophes
    return translated_text

def translate_docx(translated_file):
    doc = Document(translated_file)
    for paragraph in doc.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue  # Skip empty runs
            try:
                translated_text = translate_text(run.text)
                run.text = translated_text
            except Exception as e:
                print(f"Failed to translate text: {run.text}. Error: {str(e)}")
                continue
    
    # Translate text in tables
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        if not run.text:
                            continue  # Skip empty runs
                        try:
                            translated = translate_text(run.text)  # Translate the text to English
                            run.text = translated  # Update the text with the translated version
                        except Exception as e:
                            print(f"Failed to translate text: {run.text}. Error: {str(e)}")
                            continue

    directory, filename = os.path.split(translated_file)
    filename_without_ext = os.path.splitext(filename)[0]

    try:
        translated_filename_without_ext = translate_text(filename_without_ext)
    except Exception as e:
        print(f"Failed to translate filename: {filename_without_ext}. Error: {str(e)}")
        translated_filename_without_ext = filename_without_ext

    translated_filename = translated_filename_without_ext + ".docx"
    translated_doc_path = os.path.join(directory, translated_filename)
    doc.save(translated_doc_path)

    print(f"Translated document saved at: {translated_doc_path}")
    return translated_doc_path

def translate_text_frame(text_frame):
    for paragraph in text_frame.paragraphs:
        for run in paragraph.runs:
            if not run.text:
                continue  # Skip empty runs
            try:
                translated = translate_text(run.text)  # Translate the text to English
                run.text = translated  # Update the text with the translated version
            except Exception as e:
                print(f"Failed to translate text-frame: {run.text}. Error: {str(e)}")
                continue

def adjust_text_size(shape):
    while shape.text_frame.text != "":
        try:
            # Try to access the last character of the shape's text
            _ = shape.text_frame.text[-1]
            break  # Break the loop if the last character can be accessed
        except IndexError:
            # If the last character cannot be accessed, the text overflows the shape
            for paragraph in shape.text_frame.paragraphs:
                for run in paragraph.runs:
                    # Decrease the font size by 1 point
                    run.font.size = Pt(run.font.size.pt - 1)

def process_shape(shape):
    if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
        # If the shape is a group, recursively process each shape within the group
        for shape_in_group in shape.shapes:
            process_shape(shape_in_group)
    elif shape.has_text_frame:
        # If the shape has a text frame, translate the text and adjust the text size
        translate_text_frame(shape.text_frame)
        adjust_text_size(shape)
    elif shape.has_table:
        # If the shape is a table, process each cell's text frame and adjust text size
        for row in shape.table.rows:
            for cell in row.cells:
                translate_text_frame(cell.text_frame)
                adjust_text_size(cell)

def translate_pptx(pptx_path):
    pres = Presentation(pptx_path)

    # Process each slide in the presentation
    for slide in pres.slides:
        for shape in slide.shapes:
            process_shape(shape)
            if shape.has_text_frame:
                translate_text_frame(shape.text_frame)
            elif shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        translate_text_frame(cell.text_frame)

    directory, filename = os.path.split(pptx_path)
    filename_without_ext = os.path.splitext(filename)[0]

    try:
        translated_filename_without_ext = translate_text(filename_without_ext)
    except Exception as e:
        print(f"Failed to translate filename: {filename_without_ext}. Error: {str(e)}")
        translated_filename_without_ext = filename_without_ext

    translated_filename = translated_filename_without_ext + ".pptx"
    translated_pptx_path = os.path.join(directory, translated_filename)
    pres.save(translated_pptx_path)

    print(f"Translated presentation saved at: {translated_pptx_path}")
    return translated_pptx_path

def convert_doc_to_docx(doc_path):
    # Ensure the path is absolute
    doc_path = os.path.abspath(doc_path)

    # Create the new path by replacing the extension
    new_file_abs = doc_path.replace(".doc", ".docx")

    try:
        # Initialize the Word.Application
        word = win32com.client.Dispatch('Word.Application',pythoncom.CoInitialize())

        # Set the application to be invisible
        word.Visible = False

        # Open the document
        doc = word.Documents.Open(doc_path)

        # Save as a .docx file
        doc.SaveAs(new_file_abs, FileFormat=16)  # 16 represents the wdFormatDocx constant

        # Close the document
        doc.Close()

        # Quit Word
        word.Quit()
    except Exception as e:
        print(f"Failed to convert file: {doc_path}. Error: {str(e)}")
        return doc_path

    return new_file_abs

def main():
    source_button = QPushButton('Select Source Directory')
    target_button = QPushButton('Select Target Directory')
    start_button = QPushButton('Start Translations')
    label = QLabel()
    source_directory_label = QLabel("Source Directory: Not Selected")
    target_directory_label = QLabel("Target Directory: Not Selected")

    source_button.clicked.connect(lambda: (setattr(app, 'source_directory', select_directory()), source_directory_label.setText(f'Source Directory: {app.source_directory}'), target_button.setEnabled(True)))
    target_button.clicked.connect(lambda: (setattr(app, 'target_directory', select_directory()), target_directory_label.setText(f'Target Directory: {app.target_directory}'), start_button.setEnabled(True)))
    start_button.clicked.connect(lambda: translate_files(app.source_directory, app.target_directory, label))

    target_button.setEnabled(False)
    start_button.setEnabled(False)

    window = QWidget()
    layout = QVBoxLayout(window)
    layout.addWidget(source_button)
    layout.addWidget(source_directory_label)
    layout.addWidget(target_button)
    layout.addWidget(target_directory_label)
    layout.addWidget(start_button)
    layout.addWidget(label)
    window.show()

    app.exec_()

if __name__ == "__main__":
    main()
